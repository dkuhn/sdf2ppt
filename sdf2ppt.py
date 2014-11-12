__cmddoc__ = """

sdf2ppt - Reads an SDFile and displays molecules as image grid in powerpoint/openoffice

#Contact: daniel.kuhn@merckgroup.com - 2014 - Merck Serono
""" 

import argparse
import logging
from sys           import exit,argv
from os.path import basename,isfile

try:
  from pptx import Presentation
  from pptx.enum.text import PP_ALIGN
  from pptx.util import Inches, Px, Cm, Pt
except ImportError:
  print "Please install python-pptx module"
  exit()

try:  
  from rdkit import Chem
  from rdkit.Chem import AllChem,Draw
except:
  print "Please install/configure RDKit"
  exit()
  
  

# Path to your powerpoint template file. Change this to your needs
PPT_DEFAULT = "./my_ppt_template.pptx"



#------------------------------------------------------------------------------------------------------------------------#    
def init_logging(logfilename=None,fh_level=logging.DEBUG, ch_level=logging.INFO):
 '''
 Initialize logging
 '''
#------------------------------------------------------------------------------------------------------------------------#

 
 scriptname  = basename(argv[0].replace('.py', ''))
 # adjust name when called via python -c mechanism. Then argv == 'c'
 if scriptname == '-c':
   scriptname = 'myscript'
   #
 #
 if not logfilename:
   logfilename = './%s_%s.log' % (scriptname,timestamp())
 
 # Enable error logging via file and console
 logger = logging.getLogger("%s" % scriptname)
 #
 logger.setLevel(logging.DEBUG)
 # create file handler which logs even debug messages
 fh = logging.FileHandler(logfilename, 'w')
 fh.setLevel(fh_level)
 # create console handler with a higher log level
 ch = logging.StreamHandler()
 ch.setLevel(ch_level)
 #
 # create formatter and add it to the handlers
 #formatter = logging.Formatter("%(asctime)s - %(name)s - %(levelname)-8ss - %(message)s")
 formatter = logging.Formatter("%(asctime)s - %(name)s:  %(message)s","%d.%m.%y %H:%M:%S")
 ch.setFormatter(formatter)
 fh.setFormatter(formatter)
 #
 # add the handlers to logger
 logger.addHandler(ch)
 logger.addHandler(fh)
 return logger

#------------------------------------------------------------------------------------------------------------------------#    
def init_logging(logfilename=None,fh_level=logging.DEBUG, ch_level=logging.INFO):
 '''
 Initialize logging
 '''
#------------------------------------------------------------------------------------------------------------------------#
 from sys import argv
 from os.path import basename
 
 scriptname  = basename(argv[0].replace('.py', ''))
 # adjust name when called via python -c mechanism. Then argv == 'c'
 if scriptname == '-c':
   scriptname = 'myscript'
   #
 #
 if not logfilename:
   logfilename = './%s_%s.log' % (scriptname,timestamp())
 
 # Enable error logging via file and console
 logger = logging.getLogger("%s" % scriptname)
 #
 logger.setLevel(logging.DEBUG)
 # create file handler which logs even debug messages
 fh = logging.FileHandler(logfilename, 'w')
 fh.setLevel(fh_level)
 # create console handler with a higher log level
 ch = logging.StreamHandler()
 ch.setLevel(ch_level)
 #
 # create formatter and add it to the handlers
 #formatter = logging.Formatter("%(asctime)s - %(name)s - %(levelname)-8ss - %(message)s")
 formatter = logging.Formatter("%(asctime)s - %(name)s:  %(message)s","%d.%m.%y %H:%M:%S")
 ch.setFormatter(formatter)
 fh.setFormatter(formatter)
 #
 # add the handlers to logger
 logger.addHandler(ch)
 logger.addHandler(fh)
 return logger


def timestamp():
  import time
  lt = time.localtime(time.time())
  return "%02d_%02d_%04d__%02d_%02d_%02d" % (lt[2], lt[1], lt[0], lt[3], lt[4], lt[5])






#------------------------------------------------------------------------------------------------------------------------#    
def init_logging(logfilename=None,fh_level=logging.DEBUG, ch_level=logging.INFO):
 '''
 Initialize logging
 '''
#------------------------------------------------------------------------------------------------------------------------#
 from sys import argv
 from os.path import basename
 
 scriptname  = basename(argv[0].replace('.py', ''))
 # adjust name when called via python -c mechanism. Then argv == 'c'
 if scriptname == '-c':
   scriptname = 'myscript'
   #
 #
 if not logfilename:
   logfilename = './%s_%s.log' % (scriptname,timestamp())
 
 # Enable error logging via file and console
 logger = logging.getLogger("%s" % scriptname)
 #
 logger.setLevel(logging.DEBUG)
 # create file handler which logs even debug messages
 fh = logging.FileHandler(logfilename, 'w')
 fh.setLevel(fh_level)
 # create console handler with a higher log level
 ch = logging.StreamHandler()
 ch.setLevel(ch_level)
 #
 # create formatter and add it to the handlers
 #formatter = logging.Formatter("%(asctime)s - %(name)s - %(levelname)-8ss - %(message)s")
 formatter = logging.Formatter("%(asctime)s - %(name)s:  %(message)s","%d.%m.%y %H:%M:%S")
 ch.setFormatter(formatter)
 fh.setFormatter(formatter)
 #
 # add the handlers to logger
 logger.addHandler(ch)
 logger.addHandler(fh)
 return logger


# helper functions allowing currying of logging functionality
# http://stackoverflow.com/questions/5974273/python-avoid-passing-logger-reference-between-functions
def adjust_debug(loggername):
    logger = logging.getLogger(loggername) 
    def log_(enter_message=None, exit_message=None):
        def wrapper(f):
            #print f.func_name
            def wrapped(*args, **kargs):
                if enter_message is None:
                    logger.debug("Entering  %s" % f.func_name)
                else:
                    logger.debug(enter_message)
                r = f(*args, **kargs)
                if exit_message is None:
                    logger.debug("Leaving  %s" % f.func_name)
                else:
                    logger.debug(exit_message)
                return r
            return wrapped
        return wrapper
    return log_




# set loggername to scriptname
my_debug = adjust_debug(basename(argv[0].replace('.py', '')))




#----------------------------------------------------------------------------------------------------#
def init_current_logging():
  '''
  Wrapper functions that sets loglevel for console and file logging
  '''
#----------------------------------------------------------------------------------------------------#
  
  LOGGING_LEVELS = {'notset' : logging.NOTSET,
                    'debug': logging.DEBUG,
                    'info': logging.INFO,
                    'warning': logging.WARNING,
                    'error': logging.ERROR,
                    'critical': logging.CRITICAL}

  # Check for legal loglevels
  # set the logging_level for the console triggered via the -v switch
  ch_level= LOGGING_LEVELS['error']
  if options.verbose:
    ch_level= LOGGING_LEVELS['debug']
    # logger expects logging.DEBUG. If not in LOGGING_LEVELS dict use default parameters
    logger = init_logging(fh_level=ch_level, ch_level=ch_level)
  else:
    logger = init_logging()
    #
  #
  return logger




#----------------------------------------------------------------------------------------------------#
def init_opt_parser():
  '''
  Inits the Command-line option parser.
  '''
#----------------------------------------------------------------------------------------------------#
  
  parser = argparse.ArgumentParser(description=__cmddoc__,formatter_class=argparse.RawDescriptionHelpFormatter)
  parser.add_argument('-v','--verbose',   default=False, help="Set verbosity level", action="count")
  parser.add_argument('sdfile', action='store', default=None, help='sdfile with ligands of interest',type=str)
  parser.add_argument('--ppt_template', '-p',   action='store', default=PPT_DEFAULT, help='Powerpoint template file',type=str)
  parser.add_argument('--num_mols_page', action='store', default=8, help='Number of molecules per slide',type=int)
  
  options = parser.parse_args()
  
  
  if not isfile(options.ppt_template):
    print "Sorry, the template pptx file %s does not exist." % PPT_DEFAULT
    exit()

  return options




@my_debug()
#----------------------------------------------------------------------------------------------------#
def Mol2PNG(m,ofile='out.png',  logger=None, Sanitize=True, x=300,y=300):
  '''
  Checks Sanity of molecule and writes png
  '''
#----------------------------------------------------------------------------------------------------#
  Chem.MolToSmiles(m,True)
  if Sanitize:
    Chem.SanitizeMol(m)
  AllChem.Compute2DCoords(m)
  Draw.MolToFile(m, ofile, size=(x,y))
 
 

def sdf2ppt(mols=[]):
  mol_count=0
  page_count=1
  for mol in mols:
    # when we have num_mols_per_page mols, create new slide and header
    if mol_count % num_mols_per_page == 0:
      print "Creating new slide %i/%i" % (page_count, num_pages)
      new_slide = prs.slides.add_slide(text_slide_layout)
      title = new_slide.shapes.title      
      title.text = "Docking Results for %s  %i/%i" % (options.sdfile, page_count, num_pages)
      mol_count=0
      page_count+=1
      #
    #
    logger.debug("Processing molecule %s" % mol.GetProp("_Name"))
    Mol2PNG (mol,'%s.png' %  ( mol.GetProp("_Name")), logger=logger )
    img_path =  '%s.png' %  ( mol.GetProp("_Name"))
    #
    # track position in column 
    if mol_count % (num_mols_per_page/2) == 0:
        pos_count=0
        #
    # track position in row
    if mol_count < (num_mols_per_page/2):
      top = top_border_first_line
    else:
      top = top_border_second_line
      #
    #
    left=left_border+pos_count*image_width
    pos_count+=1
    logger.debug("Image info: %i %i %i %i" % ( mol_count, pos_count, left, top))
    #pic = new_slide.shapes.add_picture(img_path, Px(left), Px(top), Px(image_width), Px(image_height))
    pic = new_slide.shapes.add_picture(img_path, Cm(left), Cm(top), Cm(image_width), Cm(image_height))
    txBox = new_slide.shapes.add_textbox(Cm(left),  Cm(top+image_height), Cm(image_width), Cm(label_vert_size))
    tf = txBox.textframe
    tf.text = mol.GetProp("_Name")
    p = tf.paragraphs[0]
    p.font.size=Pt(12)
    p.alignment=PP_ALIGN.CENTER

    mol_count+=1
    #
  #
  prs.save(options.sdfile.replace(".sdf", ".pptx"))

   

if __name__ == '__main__':
  # Parse Commandline
  options  = init_opt_parser()
  # init logger
  logger = init_current_logging()

  
  prs = Presentation(options.ppt_template)
  text_slide_layout = prs.slide_layouts[5]

  mols  = Chem.SDMolSupplier(options.sdfile)
 

  num_mols_per_page = options.num_mols_page
  num_pages = (len(mols)/num_mols_per_page)+1
 
  # set the margins
  left_border = 1.3
  image_width = 5.3
  image_height = 5.3
  label_vert_size = 1.0
  
  # set first and second line
  top_border_first_line = 5.72
  top_border_second_line = top_border_first_line + image_width + label_vert_size
  
  sdf2ppt(mols=mols)
  

